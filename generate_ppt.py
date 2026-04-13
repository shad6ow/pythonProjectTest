from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 主题色
DARK_BLUE   = RGBColor(0x0D, 0x2B, 0x55)
ACCENT_BLUE = RGBColor(0x1F, 0x77, 0xB4)
ORANGE      = RGBColor(0xFF, 0x7F, 0x0E)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF0, 0xF4, 0xF8)
TEXT_DARK   = RGBColor(0x1A, 0x1A, 0x2E)
GREEN       = RGBColor(0x2C, 0xA0, 0x2C)
PURPLE      = RGBColor(0x8B, 0x00, 0x8B)
RED         = RGBColor(0xD6, 0x27, 0x28)
TEAL        = RGBColor(0x00, 0x86, 0x8B)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
TOTAL = 18

def rect(slide, l, t, w, h, fill=None, lc=None, lw=Pt(1)):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if lc:
        s.line.color.rgb = lc; s.line.width = lw
    else:
        s.line.fill.background()
    return s

def tb(slide, text, l, t, w, h, sz=14, bold=False,
       color=TEXT_DARK, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = True
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(sz); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = color
    return box

def header(slide, title, sub=None):
    rect(slide, 0, 0, 13.33, 1.1, fill=DARK_BLUE)
    rect(slide, 12.83, 0, 0.5, 1.1, fill=ORANGE)
    tb(slide, title, 0.3, 0.07, 12.2, 0.62, sz=22, bold=True, color=WHITE)
    if sub:
        tb(slide, sub, 0.3, 0.66, 12.2, 0.36, sz=11,
           color=RGBColor(0xAA, 0xCC, 0xEE))

def footer(slide, n):
    rect(slide, 0, 7.15, 13.33, 0.35, fill=DARK_BLUE)
    tb(slide, "联邦谱Transformer · 基于RMT的电网云边协同异常检测",
       0.3, 7.17, 9.5, 0.28, sz=9, color=RGBColor(0x88, 0xAA, 0xCC))
    tb(slide, f"{n} / {TOTAL}", 12.1, 7.17, 1.0, 0.28,
       sz=9, color=WHITE, align=PP_ALIGN.RIGHT)

def card(slide, items, l, t, w, h, title=None, tbg=ACCENT_BLUE,
         bullet="▶ ", sz=12):
    top = t
    if title:
        rect(slide, l, t, w, 0.36, fill=tbg)
        tb(slide, title, l+0.1, t+0.04, w-0.2, 0.3,
           sz=11, bold=True, color=WHITE)
        top = t + 0.36; h -= 0.36
    rect(slide, l, top, w, h, fill=LIGHT_GRAY, lc=tbg, lw=Pt(1.2))
    step = h / max(len(items), 1)
    for i, item in enumerate(items):
        tb(slide, bullet + item, l+0.12, top + 0.05 + i*step,
           w-0.25, min(step, 0.46), sz=sz, color=TEXT_DARK)

def numbox(slide, num, lab, l, t, w=2.5, h=1.0, bg=ACCENT_BLUE):
    rect(slide, l, t, w, h, fill=bg)
    tb(slide, num, l, t+0.04, w, 0.54, sz=26, bold=True,
       color=WHITE, align=PP_ALIGN.CENTER)
    tb(slide, lab, l, t+0.57, w, 0.34, sz=10,
       color=RGBColor(0xDD, 0xEE, 0xFF), align=PP_ALIGN.CENTER)

def fbox(slide, text, l, t, w, h, bg=ACCENT_BLUE, sz=11, bold=False):
    rect(slide, l, t, w, h, fill=bg)
    tb(slide, text, l+0.05, t+0.03, w-0.1, h-0.06,
       sz=sz, bold=bold, color=WHITE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════
# S1  封面
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, fill=DARK_BLUE)
rect(s, 0, 5.75, 13.33, 0.07, fill=ORANGE)
rect(s, 0, 5.82, 13.33, 1.68, fill=RGBColor(0x06, 0x1A, 0x38))
rect(s, 12.53, 0, 0.8, 5.75, fill=ACCENT_BLUE)

tb(s, "联邦谱Transformer",
   0.55, 0.85, 11.5, 1.1, sz=44, bold=True, color=WHITE)
tb(s, "基于随机矩阵理论的电网云边协同异常检测",
   0.55, 1.98, 11.0, 0.62, sz=21, color=RGBColor(0x7E, 0xC8, 0xE3))
rect(s, 0.55, 2.72, 5.8, 0.05, fill=ORANGE)
tb(s, "Federated Spectral Transformer for Smart Grid Cloud-Edge Anomaly Detection",
   0.55, 2.86, 11.5, 0.42, sz=13, color=RGBColor(0xAA, 0xCC, 0xDD))
tb(s, "数据集：SGCC 国家电网  |  42,372 用户  |  1,000+ 天",
   0.55, 3.35, 10.5, 0.38, sz=13, color=RGBColor(0x88, 0xAA, 0xBB))

for x, num, lab in [(0.55, "① RMT",  "零标签谱分析"),
                     (3.2,  "② 28维", "谱Token工程"),
                     (5.85, "③ 双路", "Transformer"),
                     (8.5,  "④ Fed",  "联邦学习"),
                     (11.1, "⑤ Stack","集成决策")]:
    rect(s, x, 4.05, 2.25, 1.25, fill=RGBColor(0x1A, 0x3A, 0x6A))
    tb(s, num, x, 4.1, 2.25, 0.58, sz=20, bold=True,
       color=ORANGE, align=PP_ALIGN.CENTER)
    tb(s, lab, x, 4.64, 2.25, 0.36, sz=10,
       color=RGBColor(0xAA, 0xCC, 0xDD), align=PP_ALIGN.CENTER)

tb(s, "中期汇报  |  汇报日期：2026 年 4 月",
   0.55, 5.42, 7.0, 0.38, sz=12, color=RGBColor(0x88, 0xAA, 0xBB))
tb(s, "集成 AUC = 0.7651  (+17.9pp vs. 零标签基线)",
   7.0, 5.42, 6.1, 0.38, sz=12, bold=True,
   color=ORANGE, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════
# S2  目录
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8, 0xFA, 0xFF))
header(s, "目  录", "Outline")
footer(s, 2)

toc = [
    ("01", "研究背景与问题定义"),
    ("02", "传统方法的局限与挑战"),
    ("03", "本文创新算法总览（五大贡献）"),
    ("04", "RMT谱分析理论与实现"),
    ("05", "28维谱Token特征工程"),
    ("06", "双路注意力Transformer"),
    ("07", "端到端完整数据流框图"),
    ("08", "联邦学习框架（FedAvg）"),
    ("09", "XGBoost集成 & Stacking"),
    ("10", "对比实验设计"),
    ("11", "消融实验设计"),
    ("12", "实验结果与性能分析"),
    ("13", "开题对应 & 后续工作"),
]
cols_toc = [0.35, 6.9]
for i, (num, title) in enumerate(toc):
    col = cols_toc[i % 2]
    y = 1.22 + (i // 2) * 0.88
    rect(s, col, y, 0.52, 0.6, fill=ACCENT_BLUE)
    tb(s, num, col, y+0.08, 0.52, 0.46, sz=13, bold=True,
       color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, col+0.52, y, 5.8, 0.6,
         fill=LIGHT_GRAY, lc=ACCENT_BLUE, lw=Pt(1))
    tb(s, title, col+0.66, y+0.12, 5.55, 0.38, sz=13, color=TEXT_DARK)

# ════════════════════════════════════════════════
# S3  研究背景与问题定义
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8, 0xFA, 0xFF))
header(s, "01  研究背景与问题定义", "Background & Problem Formulation")
footer(s, 3)

card(s, ["全球每年因窃电、设备故障造成非技术性损耗超1000亿美元",
         "中国SGCC：42,372用户，异常率约4.6%（高度不平衡）",
         "分布式新能源接入导致用电模式多变，传统阈值规则失效",
         "AMI（高级计量基础设施）产生海量时序数据，人工审核不可行"],
     0.3, 1.22, 6.1, 2.78,
     title="⚡ 电力行业背景", tbg=DARK_BLUE, sz=12)

card(s, ["数据形式：n=42,372 用户 × T≈1,000 天日均用电量时序",
         "任务：无监督/半监督识别窃电/故障用户，最大化AUC",
         "隐私约束：各变电站数据严格不可出域（电力法规）",
         "算力约束：边端设备资源受限，模型需轻量化部署"],
     6.65, 1.22, 6.4, 2.78,
     title="🔬 问题形式化定义", tbg=PURPLE, sz=12)

for x, ic, titl, desc in [
    (0.3,  "①", "标签稀缺", "正常:异常=95.4:4.6\n人工标注成本极高"),
    (4.75, "②", "数据隐私", "分布式变电站\n数据不可集中传输"),
    (9.2,  "③", "高维难题", "T=1000维时序\n传统统计方法失准"),
]:
    rect(s, x, 4.18, 3.88, 1.75, fill=RGBColor(0xE8, 0xF0, 0xFE),
         lc=ACCENT_BLUE, lw=Pt(1.5))
    tb(s, ic,   x+0.1,  4.23, 0.55, 0.55, sz=22, bold=True, color=ACCENT_BLUE)
    tb(s, titl, x+0.68, 4.23, 3.0,  0.46, sz=15, bold=True, color=DARK_BLUE)
    tb(s, desc, x+0.15, 4.78, 3.58, 1.0,  sz=12,
       color=TEXT_DARK, align=PP_ALIGN.CENTER)

rect(s, 0.3, 6.08, 12.73, 0.72, fill=DARK_BLUE)
tb(s, "研究目标：联邦谱Transformer——无监督谱分析 + 深度特征学习 + 隐私保护联邦框架，三位一体解决上述挑战",
   0.5, 6.16, 12.3, 0.55, sz=13, bold=True, color=WHITE)

# ════════════════════════════════════════════════
# S4  传统方法局限（新增核心页）
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8, 0xFA, 0xFF))
header(s, "02  传统方法的局限——为什么现有方法在电网云边场景失效",
       "Why Traditional Anomaly Detection Fails in Smart Grid Cloud-Edge Scenario")
footer(s, 4)

col_xs = [0.28, 4.08, 7.1, 10.42]
col_ws = [3.7, 2.92, 3.22, 2.68]
hdrs   = ["传统方法", "核心依赖/假设", "电网场景失效原因", "痛点"]
rect(s, 0.28, 1.2, 12.78, 0.44, fill=DARK_BLUE)
for lbl, cx, cw in zip(hdrs, col_xs, col_ws):
    tb(s, lbl, cx+0.06, 1.23, cw, 0.38, sz=11, bold=True, color=WHITE)

rows_data = [
    ("规则阈值法\n(人工设定)",
     "需要领域专家经验\n预设电量异常阈值",
     "新能源接入后用电模式持续变化\n固定阈值频繁误报漏报",
     "标签依赖\n不可泛化",
     RGBColor(0xFF, 0xEB, 0xEE)),
    ("监督学习\nSVM / RF / LSTM",
     "需要大量高质量\n异常标注样本",
     "SGCC仅4.6%异常，人工标注成本极高\n严重不平衡导致分类器退化",
     "标签依赖\n不平衡",
     RGBColor(0xFF, 0xF3, 0xE0)),
    ("孤立森林\nIsolation Forest",
     "全局i.i.d分布\n数据可集中处理",
     "隐私法规禁止数据跨站传输\n高维T=1000时孤立性度量退化",
     "隐私违规\n高维失效",
     RGBColor(0xFF, 0xFD, 0xE7)),
    ("传统PCA / 协方差\n主成分异常检测",
     "需要准确估计\n总体协方差矩阵",
     "γ=n/T≈42>>1，样本协方差严重失准\nM-P定律：噪声掩盖真实信号特征值",
     "高维小样本\n统计失准★",
     RGBColor(0xFF, 0xEB, 0xEE)),
    ("集中式深度学习\nAutoEncoder / VAE",
     "全量数据集中到\n中央服务器训练",
     "变电站数据物理隔离，集中传输\n违反合规，通信开销无法承受",
     "隐私违规\n通信瓶颈",
     RGBColor(0xFF, 0xF3, 0xE0)),
]
for ri, (mth, ass, fail, pain, bg) in enumerate(rows_data):
    y = 1.68 + ri * 1.0
    rect(s, 0.28, y, 12.78, 0.93, fill=bg, lc=ACCENT_BLUE, lw=Pt(0.5))
    tb(s, mth,  0.34,  y+0.07, 3.58, 0.8,  sz=10, bold=True, color=DARK_BLUE)
    tb(s, ass,  4.14,  y+0.07, 2.78, 0.8,  sz=10, color=TEXT_DARK)
    tb(s, fail, 7.16,  y+0.07, 3.12, 0.8,  sz=10, color=RED)
    tb(s, pain, 10.48, y+0.1,  2.5,  0.72, sz=10, bold=True,
       color=RED, align=PP_ALIGN.CENTER)

rect(s, 0.28, 6.72, 12.78, 0.38, fill=RGBColor(0xE8, 0xF5, 0xE9))
tb(s, "★ 本文核心突破：RMT自适应M-P边界精准分离信号/噪声子空间，解决高维小样本协方差失准；FedAvg解决隐私；双路Transformer解决标签稀缺",
   0.48, 6.75, 12.3, 0.3, sz=10, bold=True, color=GREEN)

# ════════════════════════════════════════════════
# S5  创新算法总览（五大贡献）
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8, 0xFA, 0xFF))
header(s, "03  本文创新算法总览——五大技术贡献",
       "Algorithm Overview: Five Key Technical Contributions")
footer(s, 5)

innov = [
    (DARK_BLUE,  "创新① 自适应RMT谱分析",
     "利用Marchenko-Pastur定律自适应判断M-P上界，通过MAD鲁棒估计σ²\n"
     "精准分离信号子空间与噪声子空间，无需任何标签\n"
     "突破：首次将随机矩阵理论用于电网用户级异常检测，AUC=0.5872（零标签基线）"),
    (ACCENT_BLUE,"创新② 28维异质谱Token工程",
     "融合4类特征：RMT谱投影(5维)+窗口统计(11维)+跨窗口增强(6维)+混沌Takens嵌入(6维)\n"
     "全向量化批量运算（无Python循环），速度提升10~50倍\n"
     "突破：将时序异常从「时间域」转化为「谱特征空间」，使深度模型可直接学习"),
    (GREEN,      "创新③ 双路注意力Transformer",
     "路径1 LocalWindowAttention(win=6)：捕捉短期用电突变（窃电操作）\n"
     "路径2 CLS+TransformerEncoder：捕捉长期趋势异常（设备老化/持续偷电）\n"
     "创新：双路融合256维特征，AdaptiveFocalLoss动态调整召回权重，AUC=0.7602"),
    (PURPLE,     "创新④ 联邦谱Transformer(FedAvg)",
     "3个变电站作为独立联邦客户端，数据严格不离站\n"
     "每轮仅传输模型参数（~200K浮点数），FedAvg加权聚合\n"
     "突破：将隐私保护深度嵌入电网检测框架，AUC与集中式训练相当"),
    (RED,        "创新⑤ 异质集成决策（Stacking）",
     "XGBoost(600树)+Transformer双路 → LogisticRegression元学习器融合\n"
     "XGBoost弥补Transformer对稀疏手工特征的劣势\n"
     "最终集成AUC=0.7651，F1=0.3317，超越所有单模型"),
]
for i, (bg, titl, desc) in enumerate(innov):
    y = 1.18 + i * 1.14
    rect(s, 0.28, y, 2.7, 0.92, fill=bg)
    tb(s, titl, 0.35, y+0.12, 2.58, 0.7, sz=11, bold=True,
       color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, 3.05, y, 9.98, 0.92, fill=LIGHT_GRAY, lc=bg, lw=Pt(1.2))
    tb(s, desc, 3.18, y+0.04, 9.75, 0.84, sz=10, color=TEXT_DARK)

rect(s, 0.28, 6.9, 12.78, 0.22, fill=ORANGE)
tb(s, "AUC 提升路径：0.5872（RMT零标签）→ 0.7602（Transformer）→ 0.7651（Stacking集成）   总提升 +17.9pp",
   0.48, 6.91, 12.3, 0.2, sz=10, bold=True, color=DARK_BLUE)

# ════════════════════════════════════════════════
# S6  RMT 谱分析理论（核心算法详解）
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8, 0xFA, 0xFF))
header(s, "04  RMT谱分析理论与实现——零标签异常检测基线",
       "Random Matrix Theory Spectral Analysis — Zero-Label Anomaly Detection")
footer(s, 6)

# 左：理论流程
rect(s, 0.28, 1.2, 5.85, 5.62, fill=LIGHT_GRAY, lc=DARK_BLUE, lw=Pt(1.5))
tb(s, "RMT算法流程", 0.35, 1.22, 5.72, 0.36, sz=13, bold=True, color=DARK_BLUE)

steps_rmt = [
    ("Step1", "滑动窗口切分",
     "window=30天, stride=7天\n → 144个窗口 (n_users, T=30)"),
    ("Step2", "MAD鲁棒估计σ²",
     "σ² = (MAD(X)/0.6745)²\n → 抗野值，高维稳健"),
    ("Step3", "计算M-P边界",
     "γ = n_sub/T,  λ± = σ²(1±√γ)²\n → 噪声特征值理论分布"),
    ("Step4", "LedoitWolf协方差",
     "收缩估计C = LedoitWolf(X_sub.T)\n → 解决高维小样本失准"),
    ("Step5", "特征值分解",
     "eigh(C) → top-k特征值/向量\n → 超出λ+即为真实信号"),
    ("Step6", "时间域投影",
     "V_time = X_sub.T @ eigvec (T×k)\n → 全量用户投影: X_all @ V_time"),
    ("Step7", "异常得分计算",
     "score_A = Σ(x·u_i)²  (信号子空间能量)\n → 用户级AUC = 0.5872"),
]
for i, (sid, sttl, sdesc) in enumerate(steps_rmt):
    y = 1.62 + i * 0.72
    rect(s, 0.35, y, 0.75, 0.55, fill=DARK_BLUE)
    tb(s, sid, 0.35, y+0.1, 0.75, 0.38, sz=9, bold=True,
       color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, sttl,  1.14, y+0.02, 1.5,  0.22, sz=10, bold=True,  color=DARK_BLUE)
    tb(s, sdesc, 1.14, y+0.24, 4.78, 0.28, sz=9,  color=TEXT_DARK)

# 右：公式说明
rect(s, 6.35, 1.2, 6.7, 5.62, fill=RGBColor(0xE8, 0xF0, 0xFE),
     lc=ACCENT_BLUE, lw=Pt(1.5))
tb(s, "核心数学公式", 6.45, 1.22, 6.55, 0.36, sz=13, bold=True, color=DARK_BLUE)

formulas = [
    ("Marchenko-Pastur 分布",
     "λ± = σ²(1 ± √γ)²,  γ = n/T\n"
     "当 λ_max > λ+ 时，存在真实信号子空间"),
    ("MAD 鲁棒方差估计",
     "σ̂² = (median|X-median(X)|/0.6745)²\n"
     "比 np.var 对野值鲁棒，适合实际电网数据"),
    ("信号子空间投影",
     "z_j = Σ_{i:λᵢ>λ+} (x_j · u_i)²\n"
     "u_i 为时间域特征向量（T维），x_j 为用户时序"),
    ("谱溢出比",
     "r = λ_max / λ+\n"
     "r>1 表示窗口内存在异常强度信号"),
    ("谱集中度",
     "c = Σᵢ λᵢ / trace(C)\n"
     "top-k特征值占总方差比例，衡量信号强度"),
]
for i, (ftitl, fform) in enumerate(formulas):
    y = 1.65 + i * 1.02
    rect(s, 6.45, y, 6.45, 0.25, fill=ACCENT_BLUE)
    tb(s, ftitl, 6.5, y+0.03, 6.35, 0.2, sz=10, bold=True, color=WHITE)
    rect(s, 6.45, y+0.25, 6.45, 0.72, fill=WHITE,
         lc=ACCENT_BLUE, lw=Pt(0.8))
    tb(s, fform, 6.55, y+0.28, 6.28, 0.62, sz=10,
       color=TEXT_DARK, italic=True)

rect(s, 0.28, 6.88, 12.78, 0.22, fill=ACCENT_BLUE)
tb(s, "关键洞察：RMT将「是否存在异常用户」转化为「协方差矩阵特征值是否超出噪声理论上界」的统计检验，完全无需标签",
   0.48, 6.9, 12.3, 0.2, sz=10, bold=True, color=WHITE)

# ════════════════════════════════════════════════
# S7  28维谱Token特征工程
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8, 0xFA, 0xFF))
header(s, "05  28维谱Token特征工程",
       "28-Dimensional Spectral Token Feature Engineering")
footer(s, 7)

layers = [
    (DARK_BLUE,  "① RMT谱投影",   "Token 1~5\n(5维)",
     "top-5特征向量投影能量 |vᵢ·x|\n逐用户×逐窗口，捕捉主成分信号方向"),
    (ACCENT_BLUE,"② 窗口统计",    "Token 6~16\n(11维)",
     "方差/均值/最大值/差分/零值率/趋势斜率/偏度\nFFT低频能量/非零比例，向量化批量计算"),
    (GREEN,      "③ 跨窗口增强",  "Token 17~21\n(5维)",
     "月度CV/季节指数/工作日比/连续零值段/滑动Z-score\n捕捉长期行为模式变化"),
    (PURPLE,     "④ 混沌延迟嵌入","Token 22~28\n(7维)",
     "Takens嵌入 d=4,τ=7，构造协方差特征值\n捕捉窃电时序的非线性混沌特征"),
]
for i, (bg, name, tokens, desc) in enumerate(layers):
    y = 1.22 + i * 1.38
    rect(s, 0.28, y, 2.3, 1.15, fill=bg)
    tb(s, name,   0.33, y+0.08, 2.22, 0.46, sz=13, bold=True,
       color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, tokens, 0.33, y+0.6,  2.22, 0.42, sz=11,
       color=RGBColor(0xDD, 0xEE, 0xFF), align=PP_ALIGN.CENTER)
    rect(s, 2.65, y+0.12, 8.58, 0.9, fill=LIGHT_GRAY, lc=bg, lw=Pt(1.2))
    tb(s, desc, 2.78, y+0.2, 8.35, 0.72, sz=11, color=TEXT_DARK)
    rect(s, 11.3, y+0.12, 1.75, 0.9, fill=bg)
    tb(s, f"AUC贡献\n{['强','强','中','中'][i]}",
       11.3, y+0.22, 1.75, 0.65, sz=10, bold=True,
       color=WHITE, align=PP_ALIGN.CENTER)

# 归一化策略
rect(s, 0.28, 6.72, 6.1, 0.38, fill=DARK_BLUE)
tb(s, "归一化策略：Token1~5 → 逐用户MinMax；Token6~28 → 全局RobustScaler(IQR)",
   0.38, 6.76, 5.9, 0.28, sz=9, color=WHITE)
rect(s, 6.45, 6.72, 6.6, 0.38, fill=GREEN)
tb(s, "⚡ 全向量化：消除for循环，提取速度提升10~50×",
   6.55, 6.76, 6.4, 0.28, sz=9, bold=True, color=WHITE)

# ════════════════════════════════════════════════
# S8  双路 Transformer 架构框图
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8, 0xFA, 0xFF))
header(s, "06  双路注意力Transformer——模型架构框图",
       "Dual-Path Attention Transformer Architecture")
footer(s, 8)

# 输入层
fbox(s, "输入\n(B, 48, 28)\n48窗口×28维Token",
     0.28, 1.2, 2.25, 1.3, bg=DARK_BLUE, sz=11)
tb(s, "→", 2.55, 1.6, 0.35, 0.5,
   sz=18, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

# 路径1 LocalWindow
rect(s, 2.95, 1.2, 4.5, 5.52, fill=RGBColor(0xE3, 0xF2, 0xFD),
     lc=ACCENT_BLUE, lw=Pt(2))
tb(s, "路径1: 局部窗口注意力", 3.0, 1.22, 4.4, 0.36,
   sz=11, bold=True, color=ACCENT_BLUE)
path1_steps = [
    ("Linear Proj", "d_model=128"),
    ("+ PositionEmb", "正弦位置编码"),
    ("LocalWinAttn×4", "win=6, Pre-LN"),
    ("FFN×4",         "dim_ff=512"),
    ("CLS Pooling",   "→(B,128) feat_local"),
]
for i, (nm, ds) in enumerate(path1_steps):
    y = 1.68 + i * 0.95
    fbox(s, f"{nm}\n{ds}", 3.1, y, 4.2, 0.75,
         bg=ACCENT_BLUE, sz=10)

# 路径2 Global CLS
rect(s, 7.62, 1.2, 4.5, 5.52, fill=RGBColor(0xF3, 0xE5, 0xF5),
     lc=PURPLE, lw=Pt(2))
tb(s, "路径2: 全局CLS Transformer", 7.67, 1.22, 4.4, 0.36,
   sz=11, bold=True, color=PURPLE)
path2_steps = [
    ("Linear Proj", "d_model=128"),
    ("CLS Token插入", "[CLS] + seq"),
    ("TransformerEnc×4", "4heads, Pre-LN"),
    ("FFN×4",             "dim_ff=512"),
    ("CLS Pooling",       "→(B,128) feat_global"),
]
for i, (nm, ds) in enumerate(path2_steps):
    y = 1.68 + i * 0.95
    fbox(s, f"{nm}\n{ds}", 7.72, y, 4.2, 0.75, bg=PURPLE, sz=10)

# 融合层
tb(s, "↓ concat\n(B,256)", 2.0, 3.0, 0.9, 0.7,
   sz=9, color=ORANGE, align=PP_ALIGN.CENTER, bold=True)

fbox(s, "特征融合 Concat\nfeat_local(B,128) + feat_global(B,128)\n= feat_cat(B,256)",
     3.1, 6.82, 8.42, 0.2, bg=RGBColor(0x3E, 0x2D, 0x70), sz=9)

# 输出
rect(s, 12.15, 1.2, 0.9, 5.52, fill=RGBColor(0xE8, 0xF5, 0xE9),
     lc=GREEN, lw=Pt(2))
tb(s, "输\n出\n层", 12.18, 1.3, 0.82, 1.2, sz=11, bold=True,
   color=GREEN, align=PP_ALIGN.CENTER)
out_steps = ["UserGraph\nAttn(ep30+)",
             "Dropout\n(p=0.2)",
             "MLP\n(256→64→1)",
             "Sigmoid\n→ logit",
             "AdaptiveFocal\nLoss"]
for i, st in enumerate(out_steps):
    fbox(s, st, 12.18, 2.6+i*0.82, 0.82, 0.7, bg=GREEN, sz=8)

# ════════════════════════════════════════════════
# S9  端到端完整数据流框图
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8, 0xFA, 0xFF))
header(s, "07  端到端完整数据流框图",
       "End-to-End Data Flow Diagram")
footer(s, 9)

# 六步流程 主干
steps6 = [
    (DARK_BLUE,  "① 数据预处理",
     "原始SGCC CSV\n↓线性插值\n↓Z-score归一化\n↓极端值裁剪\n(42372, 1000)"),
    (ACCENT_BLUE,"② RMT谱分析",
     "滑动窗口×144\n↓LedoitWolf\n↓M-P边界\n↓特征值分解\n→谱Token7维"),
    (GREEN,      "③ 28维特征工程",
     "Token1~5 谱投影\n↓+Token6~16统计\n↓+Token17~21跨窗口\n↓+Token22~28混沌\n→(42372,48,28)"),
    (PURPLE,     "④ 双路Transformer",
     "LocalWinAttn×4\n+CLS Global×4\n↓feat_cat(B,256)\n↓AdaptiveFocalLoss\n→AUC=0.7602"),
    (ORANGE,     "⑤ FedAvg聚合",
     "3站点本地训练\n↓上传∆W\n↓加权聚合\nw=Σ(nk/n)wk\n→隐私保护"),
    (RED,        "⑥ Stacking集成",
     "Transformer概率\n+XGBoost概率\n+256维深度特征\n↓LogReg融合\n→AUC=0.7651"),
]
bw = 1.82
for i, (bg, ttl, desc) in enumerate(steps6):
    x = 0.22 + i * 2.18
    rect(s, x, 1.3, bw, 0.42, fill=bg)
    tb(s, ttl, x, 1.32, bw, 0.38, sz=11, bold=True,
       color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x, 1.72, bw, 2.8, fill=LIGHT_GRAY, lc=bg, lw=Pt(1.5))
    tb(s, desc, x+0.06, 1.8, bw-0.12, 2.65, sz=10,
       color=TEXT_DARK, align=PP_ALIGN.CENTER)
    if i < 5:
        tb(s, "→", x + bw, 2.72, 0.36, 0.5,
           sz=20, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

# 三大子系统说明
sys_info = [
    (0.22, DARK_BLUE,  "无监督谱分析子系统（边端执行）",
     "RMT+特征工程：无需标签，计算轻量\n可在变电站边端独立运行"),
    (4.58, ACCENT_BLUE,"深度学习子系统（本地训练）",
     "双路Transformer：深度挖掘时序异常模式\n参数量~200K，支持边端部署"),
    (8.94, RED,        "联邦聚合子系统（云端协调）",
     "FedAvg+Stacking：隐私保护聚合\n最终决策在云端，原始数据不出站"),
]
for x, bg, titl, desc in sys_info:
    rect(s, x, 4.62, 4.14, 1.48, fill=RGBColor(0xF8, 0xFA, 0xFF),
         lc=bg, lw=Pt(2))
    tb(s, titl, x+0.1, 4.65, 3.95, 0.38, sz=11, bold=True, color=bg)
    tb(s, desc, x+0.1, 5.08, 3.95, 0.95, sz=10, color=TEXT_DARK)

rect(s, 0.22, 6.18, 12.89, 0.72, fill=DARK_BLUE)
tb(s, "✅ 全流程已完整实现：RMT(0.5872) → Transformer(0.7602) → Stacking(0.7651)  累计提升 +17.9pp",
   0.42, 6.28, 12.48, 0.52, sz=13, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════
# S10  联邦学习框架
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8, 0xFA, 0xFF))
header(s, "08  联邦学习框架（FedAvg）——隐私保护协同训练",
       "Federated Learning Framework — Privacy-Preserving Collaborative Training")
footer(s, 10)

# 三个客户端
for i, (lab, users) in enumerate([("变电站 A", "~14,124 户"),
                                    ("变电站 B", "~14,124 户"),
                                    ("变电站 C", "~14,124 户")]):
    x = 0.3 + i * 3.82
    rect(s, x, 1.25, 3.4, 2.38, fill=ACCENT_BLUE)
    tb(s, lab,   x, 1.33, 3.4, 0.52, sz=18, bold=True,
       color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, users, x, 1.84, 3.4, 0.38, sz=13,
       color=RGBColor(0xDD, 0xEE, 0xFF), align=PP_ALIGN.CENTER)
    rect(s, x+0.1, 2.26, 3.2, 1.2, fill=RGBColor(0x15, 0x5E, 0xA8))
    tb(s, "本地训练\n1轮 × local_epochs=3\n上传 ΔW（~200K参数）",
       x+0.15, 2.32, 3.1, 1.08, sz=10, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, "↑ 上传ΔW   ↓ 广播W_global",
       x+0.3, 3.68, 2.85, 0.38, sz=10, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

# 云端聚合
rect(s, 4.2, 3.95, 4.93, 1.58, fill=DARK_BLUE)
tb(s, "☁  云端 FedAvg 聚合",
   4.2, 4.03, 4.93, 0.48, sz=15, bold=True,
   color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "w_global = Σₖ (nₖ / n_total) · wₖ",
   4.2, 4.55, 4.93, 0.42, sz=12, color=ORANGE, align=PP_ALIGN.CENTER)
tb(s, "预创建本地模型（避免重复__init__）",
   4.2, 5.0, 4.93, 0.42, sz=10, color=RGBColor(0xAA, 0xCC, 0xDD),
   align=PP_ALIGN.CENTER)

card(s, ["🔒 隐私保护：原始用电数据严格不离开本地变电站",
         "📡 通信效率：每轮仅传输~200K浮点参数（非原始数据）",
         "⚡ 工程优化：预创建local_models复用，torch.compile加速",
         "✅ 实验验证：联邦AUC与集中式训练相当，隐私代价极小"],
     0.3, 5.65, 12.73, 1.38,
     title="联邦学习关键结论", tbg=GREEN, sz=12)

# ════════════════════════════════════════════════
# S11  XGBoost 集成
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8, 0xFA, 0xFF))
header(s, "09  XGBoost集成 & Stacking——异质模型融合",
       "Ensemble Learning: XGBoost + Stacking Meta-Learner")
footer(s, 11)

card(s, ["提取双路Transformer 256维中间特征",
         "  feat_local(B,128) + feat_global(B,128)",
         "拼接224维手工统计特征（均值/方差/分位数）",
         "XGBoost：600棵树，max_depth=6，early stopping",
         "XGBoost 单独 AUC = 0.7560",
         "Stacking：LogisticRegression 融合两路输出概率"],
     0.28, 1.22, 6.2, 4.55,
     title="🔧 实现细节", tbg=DARK_BLUE, sz=12)

# Stacking流程图
rect(s, 6.65, 1.22, 6.4, 4.55, fill=RGBColor(0xE8, 0xF0, 0xFE),
     lc=ACCENT_BLUE, lw=Pt(1.5))
tb(s, "Stacking 融合流程", 6.75, 1.24, 6.2, 0.36,
   sz=12, bold=True, color=DARK_BLUE)

fbox(s, "Transformer 预测概率 p₁",    6.85, 1.68, 5.95, 0.55, bg=PURPLE, sz=11)
fbox(s, "XGBoost 预测概率 p₂",        6.85, 2.32, 5.95, 0.55, bg=ACCENT_BLUE, sz=11)
fbox(s, "Transformer 256维深度特征 f", 6.85, 2.96, 5.95, 0.55, bg=GREEN, sz=11)
tb(s, "↓  LogisticRegression 元学习器  ↓",
   6.85, 3.62, 5.95, 0.42, sz=11, bold=True,
   color=DARK_BLUE, align=PP_ALIGN.CENTER)
fbox(s, "最终异常评分  AUC = 0.7651\nF1 = 0.3317",
     6.85, 4.12, 5.95, 0.68, bg=RED, sz=13, bold=True)

rect(s, 0.28, 5.88, 12.78, 0.68, fill=RGBColor(0xFF, 0xF3, 0xE0))
tb(s, "💡 集成原理：XGBoost 擅长稀疏手工统计特征，Transformer 擅长时序深度模式\n"
   "两者互补性经 Stacking 融合后实现超越单模型的性能突破",
   0.45, 5.93, 12.4, 0.58, sz=11, color=RGBColor(0x8B, 0x45, 0x00))

# ════════════════════════════════════════════════
# S12  对比实验设计
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8, 0xFA, 0xFF))
header(s, "10  对比实验设计",
       "Comparative Experiment Design")
footer(s, 12)

rect(s, 0.28, 1.2, 12.78, 0.42, fill=DARK_BLUE)
for lbl, cx, cw in zip(["对比方法", "方法类别", "核心配置", "预期/已有AUC"],
                        [0.3, 3.4, 6.0, 11.2],
                        [3.0, 2.5, 5.1, 1.8]):
    tb(s, lbl, cx+0.06, 1.23, cw, 0.36, sz=11, bold=True, color=WHITE)

comp_rows = [
    ("RMT谱分析（本文）",    "无监督基线",  "M-P边界+LedoitWolf+MAD估计σ²",           "0.5872 ✅已复现", RGBColor(0xE3, 0xF2, 0xFD), True),
    ("双路Transformer（本文）","深度学习",   "LocalWin+CLS, d=128, 4层, AdaptiveFocal", "0.7602 ✅已复现", RGBColor(0xE3, 0xF2, 0xFD), True),
    ("本文集成（本文）",       "集成模型",   "Transformer+XGBoost+Stacking",            "0.7651 ✅已复现", RGBColor(0xE8, 0xF5, 0xE9), True),
    ("LSTM-AE",               "深度无监督", "2层LSTM编码器，重构误差为异常分",           "待实验",          LIGHT_GRAY, False),
    ("VAE（变分自编码器）",    "深度无监督", "潜变量采样，ELBO损失",                     "待实验",          LIGHT_GRAY, False),
    ("One-Class SVM",          "传统无监督", "RBF核，nu=0.05",                          "待实验",          LIGHT_GRAY, False),
    ("孤立森林 IF",            "传统无监督", "n_estimators=200，contamination=0.046",    "待实验",          LIGHT_GRAY, False),
    ("FedProx（联邦变体）",   "联邦学习",   "近端正则项μ=0.01，对比FedAvg隐私效用",    "待实验",          LIGHT_GRAY, False),
]
for ri, (nm, cat, cfg, auc, bg, done) in enumerate(comp_rows):
    y = 1.66 + ri * 0.68
    rect(s, 0.28, y, 12.78, 0.62, fill=bg, lc=ACCENT_BLUE, lw=Pt(0.4))
    fc = GREEN if done else TEXT_DARK
    tb(s, nm,  0.34, y+0.07, 2.98, 0.5, sz=10, bold=done, color=fc)
    tb(s, cat, 3.46, y+0.07, 2.42, 0.5, sz=10, color=TEXT_DARK)
    tb(s, cfg, 6.06, y+0.07, 5.02, 0.5, sz=9,  color=TEXT_DARK)
    tb(s, auc, 11.26, y+0.07, 1.7, 0.5, sz=10, bold=done,
       color=GREEN if done else RGBColor(0xAA, 0xAA, 0xAA),
       align=PP_ALIGN.CENTER)

rect(s, 0.28, 7.08, 12.78, 0.02, fill=ACCENT_BLUE)

# ════════════════════════════════════════════════
# S13  消融实验设计
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8, 0xFA, 0xFF))
header(s, "11  消融实验设计——验证各组件贡献",
       "Ablation Study Design — Validating Each Component's Contribution")
footer(s, 13)

card(s, ["去除RMT谱Token(Token1~5)，仅保留统计特征",
         "→ 验证RMT谱投影对最终AUC的独立贡献",
         "假设：AUC显著下降（预计-3~5pp）"],
     0.28, 1.22, 6.25, 1.78,
     title="消融① RMT谱特征的贡献", tbg=DARK_BLUE, sz=11)

card(s, ["去除Takens混沌嵌入(Token22~28)",
         "→ 验证非线性特征对窃电检测的增益",
         "假设：对复杂窃电模式AUC下降约2pp"],
     6.65, 1.22, 6.4, 1.78,
     title="消融② 混沌Takens嵌入的贡献", tbg=PURPLE, sz=11)

card(s, ["仅用单路LocalWindowAttention（去除CLS全局路径）",
         "→ 验证双路融合vs单路的性能差异",
         "假设：缺少全局依赖导致AUC下降约2~4pp"],
     0.28, 3.1, 6.25, 1.78,
     title="消融③ 双路架构 vs 单路架构", tbg=ACCENT_BLUE, sz=11)

card(s, ["去除UserGraphAttention（批内用户图注意力）",
         "→ 验证图结构用户关系建模的贡献",
         "假设：批内异常传播信号损失，AUC下降约1~2pp"],
     6.65, 3.1, 6.4, 1.78,
     title="消融④ UserGraphAttention的贡献", tbg=GREEN, sz=11)

card(s, ["固定AdaptiveFocalLoss权重=1.0（去除自适应机制）",
         "→ 验证动态召回权重对不平衡数据的效果",
         "假设：F1显著下降（AUC影响相对较小）"],
     0.28, 4.98, 6.25, 1.78,
     title="消融⑤ AdaptiveFocalLoss的贡献", tbg=RED, sz=11)

card(s, ["去除Stacking，仅用单一Transformer输出",
         "→ 验证异质集成相对于单模型的增益",
         "假设：AUC下降约0.5~1.5pp（验证集成必要性）"],
     6.65, 4.98, 6.4, 1.78,
     title="消融⑥ Stacking集成的必要性", tbg=ORANGE, sz=11)

rect(s, 0.28, 6.84, 12.78, 0.26, fill=DARK_BLUE)
tb(s, "消融实验策略：每次仅去除/替换一个组件，其余保持最优配置，5折交叉验证取均值，确保统计显著性",
   0.42, 6.86, 12.3, 0.22, sz=10, color=WHITE)

# ════════════════════════════════════════════════
# S14  实验结果与性能分析
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8, 0xFA, 0xFF))
header(s, "12  实验结果与性能分析",
       "Experimental Results & Performance Analysis")
footer(s, 14)

# 性能表格
hdrs_t = ["方  法", "AUC", "F1", "说  明"]
col_xt = [0.28, 5.2, 7.55, 9.0]
col_wt = [4.82, 2.25, 1.35, 4.1]
rect(s, 0.28, 1.2, 12.78, 0.48, fill=DARK_BLUE)
for hdr, cx, cw in zip(hdrs_t, col_xt, col_wt):
    tb(s, hdr, cx+0.06, 1.24, cw, 0.4, sz=13, bold=True, color=WHITE)

result_rows = [
    ("RMT谱得分（零标签基线）",      "0.5872", "N/A",   "无需任何标签，纯随机矩阵理论",
     RGBColor(0xF5,0xF5,0xF5), False),
    ("双路Transformer（28维Token）",  "0.7602", "0.3410","LocalWin+CLS双路，AdaptiveFocal",
     LIGHT_GRAY, False),
    ("XGBoost（Transformer深度特征）","0.7560", "N/A",   "256+224维特征，600棵树",
     RGBColor(0xF5,0xF5,0xF5), False),
    ("联邦Transformer（FedAvg）",     "~0.760", "~0.34", "3客户端联邦训练，与集中式相当",
     LIGHT_GRAY, False),
    ("集成（Transformer+XGBoost）",   "0.7651", "0.3317","Stacking元学习融合，最优",
     RGBColor(0xE8,0xF5,0xE9), True),
]
for ri, (nm, auc, f1, note, bg, best) in enumerate(result_rows):
    y = 1.72 + ri * 0.95
    rect(s, 0.28, y, 12.78, 0.88, fill=bg, lc=ACCENT_BLUE, lw=Pt(0.5))
    for val, cx, cw in zip([nm, auc, f1, note], col_xt, col_wt):
        fc = GREEN if best else TEXT_DARK
        bd = best and val in (auc, f1)
        tb(s, val, cx+0.06, y+0.1, cw, 0.7,
           sz=14 if bd else 12, bold=bd, color=fc)
    if best:
        tb(s, "🏆", 12.85, y+0.1, 0.4, 0.7, sz=16)

# AUC提升条
rect(s, 0.28, 6.5, 12.78, 0.65, fill=DARK_BLUE)
tb(s, "📈 AUC提升路径：0.5872（RMT零标签）→ 0.7602（Transformer）→ 0.7651（集成）  总提升 +17.9pp",
   0.48, 6.58, 12.3, 0.48, sz=13, bold=True,
   color=ORANGE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════
# S15  开题对应 & 后续工作
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8, 0xFA, 0xFF))
header(s, "13  开题对应关系 & 后续工作计划",
       "Alignment with Proposal & Future Work")
footer(s, 15)

card(s, ["✅ ① 零标签学习范式（RMT）— 自适应M-P边界+MAD估计已实现",
         "✅ ② 谱特征+非线性融合 — 28维Token+双路Transformer完成",
         "✅ ③ 联邦协同机制 — FedAvg 3客户端联邦训练验证完毕",
         "✅ ④ 工程优化 — 向量化特征提取+torch.compile+预创建模型",
         "✅ ⑤ 实验闭环 — RMT→特征→Transformer→Stacking全链路"],
     0.28, 1.22, 6.25, 4.52,
     title="📋 已完成（对应开题）", tbg=GREEN, sz=12)

card(s, ["🔄 SA/PSO启发式搜索最优RMT协方差子空间",
         "🔄 PMU/真实电网数据（非AMI）验证泛化性",
         "🔄 完成LSTM-AE/VAE/OCSVM对比实验",
         "🔄 FedProx vs FedAvg通信效率量化对比",
         "🔄 完成全部6项消融实验并统计显著性",
         "🔄 理论推导：RMT高维小样本误判率上界",
         "🔄 撰写论文（目标：IEEE Trans on Smart Grid）"],
     6.65, 1.22, 6.4, 4.52,
     title="🚀 下一步工作（3个月计划）", tbg=RED, sz=12)

rect(s, 0.28, 5.88, 12.78, 0.88, fill=DARK_BLUE)
tb(s, "💡 核心结论：本文提出的联邦谱Transformer框架将三大挑战（标签稀缺/数据隐私/高维难题）一体化解决\n"
   "RMT理论与深度学习的有机融合是本文最核心的创新，AUC已从0.5872提升至0.7651（+17.9pp）",
   0.45, 5.94, 12.3, 0.75, sz=12, bold=True, color=WHITE)

# ════════════════════════════════════════════════
# S16  致谢
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, fill=DARK_BLUE)
rect(s, 0, 5.75, 13.33, 0.07, fill=ORANGE)
rect(s, 0, 5.82, 13.33, 1.68, fill=RGBColor(0x06, 0x1A, 0x38))
rect(s, 12.53, 0, 0.8, 5.75, fill=ACCENT_BLUE)

tb(s, "感谢聆听", 1.0, 1.5, 11.0, 1.25,
   sz=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
rect(s, 3.5, 2.95, 6.33, 0.06, fill=ORANGE)
tb(s, "Thank You & Questions Welcome",
   1.0, 3.1, 11.0, 0.58, sz=20,
   color=RGBColor(0x7E, 0xC8, 0xE3), align=PP_ALIGN.CENTER)

for x, num, lab, bg in [(1.5,  "0.5872", "RMT 基线 AUC",     ACCENT_BLUE),
                         (5.4,  "0.7602", "Transformer AUC",  GREEN),
                         (9.3,  "0.7651", "集成最终 AUC",      RED)]:
    rect(s, x, 3.88, 2.5, 1.28, fill=bg)
    tb(s, num, x, 3.95, 2.5, 0.62, sz=26, bold=True,
       color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, lab, x, 4.55, 2.5, 0.42, sz=11,
       color=RGBColor(0xDD, 0xEE, 0xFF), align=PP_ALIGN.CENTER)

tb(s, "联邦谱Transformer：基于随机矩阵理论的电网云边协同异常检测",
   1.0, 6.05, 11.0, 0.45, sz=13,
   color=RGBColor(0x88, 0xAA, 0xBB), align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════
# 保存
# ════════════════════════════════════════════════
output = "联邦谱Transformer_中期汇报.pptx"
prs.save(output)
import sys
sys.stdout.buffer.write(f"PPT OK: {output}, slides={len(prs.slides)}\n".encode('utf-8'))
